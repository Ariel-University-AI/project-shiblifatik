from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colors ──────────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x0d, 0x1f, 0x38)
MID_BLUE    = RGBColor(0x1a, 0x4a, 0x7a)
ACCENT_BLUE = RGBColor(0x2d, 0x6a, 0x9f)
LIGHT_BLUE  = RGBColor(0xcf, 0xe2, 0xf7)
WHITE       = RGBColor(0xff, 0xff, 0xff)
ORANGE      = RGBColor(0xe6, 0x51, 0x00)
GREEN       = RGBColor(0x1b, 0x5e, 0x20)
GOLD        = RGBColor(0xf5, 0xa6, 0x23)

W = Inches(13.33)
H = Inches(7.5)


def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def bg(slide, color=DARK_BLUE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=24, bold=False, color=WHITE, align=PP_ALIGN.RIGHT, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def accent_bar(slide, top, color=ACCENT_BLUE):
    box(slide, Inches(0), top, Inches(0.08), Inches(0.7), color)


# ════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Slide 1: Title ──────────────────────────────────────────────────────────
s = add_slide(prs)
bg(s, DARK_BLUE)

# Gradient bar top
box(s, 0, 0, W, Inches(0.12), ACCENT_BLUE)

# Main title
txt(s, "📐 Antigrafity", Inches(1), Inches(1.5), Inches(11), Inches(1.4),
    size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
txt(s, "מערכת ניהול חכמה למשרד מדידות קרקע",
    Inches(1), Inches(3.0), Inches(11), Inches(0.9),
    size=26, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# Tags row
tags = ["Python", "Streamlit", "Pandas", "Plotly", "scikit-learn"]
tag_w = Inches(1.7)
start = Inches(2.6)
for i, tag in enumerate(tags):
    bx = box(s, start + i * (tag_w + Inches(0.15)),
             Inches(4.3), tag_w, Inches(0.5), ACCENT_BLUE)
    txt(s, tag,
        start + i * (tag_w + Inches(0.15)) + Inches(0.05),
        Inches(4.32), tag_w - Inches(0.1), Inches(0.5),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Bottom bar
box(s, 0, H - Inches(0.12), W, Inches(0.12), MID_BLUE)

# ── Slide 2: הבעיה ──────────────────────────────────────────────────────────
s = add_slide(prs)
bg(s, DARK_BLUE)
box(s, 0, 0, W, Inches(0.12), ACCENT_BLUE)

txt(s, "הבעיה", Inches(0.5), Inches(0.25), Inches(6), Inches(0.8),
    size=36, bold=True, color=WHITE)

problems = [
    ("📊", "ניהול ב-Excel בלבד",     "מאות פרויקטים, אין ניראות, אין סינון מהיר"),
    ("🔍", "כפילויות גוש/חלקה",      "פרויקט נפתח מחדש על אותה חלקה — לא מזוהה"),
    ("⏱️", "אין מעקב שעות",          "אי אפשר לדעת כמה זמן כל עבודה לוקחת"),
    ("📋", "אין ניהול משימות",        "מה עושה כל עובד? מה תקוע? לא ידוע"),
]

for i, (icon, title, desc) in enumerate(problems):
    top = Inches(1.4 + i * 1.35)
    box(s, Inches(0.4), top, Inches(12.1), Inches(1.1), MID_BLUE)
    txt(s, icon + "  " + title, Inches(0.6), top + Inches(0.08),
        Inches(4), Inches(0.55), size=20, bold=True, color=WHITE)
    txt(s, desc, Inches(4.8), top + Inches(0.08),
        Inches(7.5), Inches(0.55), size=18, color=LIGHT_BLUE)

# ── Slide 3: הפתרון ─────────────────────────────────────────────────────────
s = add_slide(prs)
bg(s, DARK_BLUE)
box(s, 0, 0, W, Inches(0.12), ACCENT_BLUE)

txt(s, "הפתרון — 6 מודולים", Inches(0.5), Inches(0.25),
    Inches(10), Inches(0.8), size=36, bold=True, color=WHITE)

modules = [
    ("📈", "EDA",           "ניתוח 310 עבודות היסטוריות"),
    ("🤖", "AI / ML",       "חיזוי זמן מדידה — scikit-learn"),
    ("🏠", "דשבורד",        "KPIs: פרויקטים, משימות, דד-ליין"),
    ("📁", "פרויקטים",      "CRUD + זיהוי כפילויות אוטומטי"),
    ("✅", "משימות",        "Kanban + שעות + ייצוא Excel"),
    ("📖", "הסברי גרפים",   "תיעוד אינטראקטיבי"),
]

col_w = Inches(4.1)
for i, (icon, title, desc) in enumerate(modules):
    col = i % 3
    row = i // 3
    left = Inches(0.4) + col * (col_w + Inches(0.25))
    top  = Inches(1.35) + row * Inches(2.6)
    bx   = box(s, left, top, col_w, Inches(2.3), MID_BLUE)
    txt(s, icon, left + Inches(0.2), top + Inches(0.2),
        Inches(0.8), Inches(0.7), size=30, align=PP_ALIGN.LEFT)
    txt(s, title, left + Inches(0.2), top + Inches(0.85),
        col_w - Inches(0.4), Inches(0.55),
        size=20, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    txt(s, desc, left + Inches(0.2), top + Inches(1.45),
        col_w - Inches(0.4), Inches(0.6),
        size=14, color=LIGHT_BLUE, align=PP_ALIGN.RIGHT)

# ── Slide 4: נתוני מפתח ────────────────────────────────────────────────────
s = add_slide(prs)
bg(s, DARK_BLUE)
box(s, 0, 0, W, Inches(0.12), ACCENT_BLUE)

txt(s, "נתוני מפתח", Inches(0.5), Inches(0.25),
    Inches(8), Inches(0.8), size=36, bold=True, color=WHITE)

stats = [
    ("310", "פרויקטים בנתוני האימון"),
    ("6",   "עמודי ממשק"),
    ("0",   "תלות שרת — הכל מקומי"),
    ("3",   "עובדים מנוהלים"),
]

stat_w = Inches(2.9)
for i, (num, label) in enumerate(stats):
    left = Inches(0.4) + i * (stat_w + Inches(0.3))
    bx = box(s, left, Inches(1.4), stat_w, Inches(2.2), ACCENT_BLUE)
    txt(s, num, left, Inches(1.55), stat_w, Inches(1.1),
        size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, label, left + Inches(0.1), Inches(2.75),
        stat_w - Inches(0.2), Inches(0.7),
        size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# Features list
features = [
    "✅  RTL עברית מלאה",
    "✅  מודל ML מאומן על נתוני אמת",
    "✅  ייצוא Excel ישיר מכל עמוד",
    "✅  זיהוי אוטומטי כפילויות גוש/חלקה",
    "✅  דשבורד KPI בזמן אמת",
]
for i, feat in enumerate(features):
    txt(s, feat, Inches(0.5), Inches(4.0) + i * Inches(0.58),
        Inches(12), Inches(0.5), size=17, color=WHITE)

# ── Slide 5: ארכיטקטורה ─────────────────────────────────────────────────────
s = add_slide(prs)
bg(s, DARK_BLUE)
box(s, 0, 0, W, Inches(0.12), ACCENT_BLUE)

txt(s, "ארכיטקטורה טכנית", Inches(0.5), Inches(0.25),
    Inches(10), Inches(0.8), size=36, bold=True, color=WHITE)

layers = [
    (WHITE,       "ממשק משתמש",  "Streamlit + Plotly + HTML/CSS RTL"),
    (LIGHT_BLUE,  "לוגיקה",      "Python — pandas, pathlib, datetime"),
    (ACCENT_BLUE, "ML",          "scikit-learn — regression model"),
    (MID_BLUE,    "נתונים",      "CSV files — projects, tasks, 2025_cleaned"),
]

layer_colors = [
    RGBColor(0x1e, 0x3a, 0x5f),
    RGBColor(0x15, 0x4a, 0x7a),
    RGBColor(0x0d, 0x3b, 0x6e),
    RGBColor(0x0d, 0x1f, 0x38),
]

for i, (_, title, desc) in enumerate(layers):
    top = Inches(1.35) + i * Inches(1.35)
    bx = box(s, Inches(0.4), top, Inches(12.1), Inches(1.15),
             layer_colors[i])
    box(s, Inches(0.4), top, Inches(0.08), Inches(1.15), GOLD)
    txt(s, title, Inches(0.7), top + Inches(0.1),
        Inches(3.5), Inches(0.55), size=20, bold=True, color=GOLD)
    txt(s, desc, Inches(4.5), top + Inches(0.1),
        Inches(7.8), Inches(0.55), size=18, color=LIGHT_BLUE)

# ── Slide 6: Demo Flow ───────────────────────────────────────────────────────
s = add_slide(prs)
bg(s, DARK_BLUE)
box(s, 0, 0, W, Inches(0.12), ACCENT_BLUE)

txt(s, "מהלך הדמו", Inches(0.5), Inches(0.25),
    Inches(8), Inches(0.8), size=36, bold=True, color=WHITE)

steps = [
    ("1", "דשבורד ראשי",       "הצגת KPIs — פרויקטים פתוחים, משימות תקועות, דד-ליין"),
    ("2", "ניתוח נתונים",      "טעינת 310 עבודות, גרפים, פילטרים"),
    ("3", "מודל AI",           "הזן ערכים → קבל חיזוי זמן מדידה"),
    ("4", "ניהול פרויקטים",   "הוסף פרויקט → זיהוי כפילות אוטומטי"),
    ("5", "משימות",            "Kanban → שנה סטטוס → ייצא Excel"),
]

for i, (num, title, desc) in enumerate(steps):
    top = Inches(1.3) + i * Inches(1.15)
    # Number circle
    box(s, Inches(0.4), top + Inches(0.05), Inches(0.7), Inches(0.7), ACCENT_BLUE)
    txt(s, num, Inches(0.4), top + Inches(0.02),
        Inches(0.7), Inches(0.7), size=22, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, Inches(1.3), top + Inches(0.05),
        Inches(3.2), Inches(0.6), size=20, bold=True, color=WHITE)
    txt(s, desc, Inches(4.7), top + Inches(0.05),
        Inches(8.0), Inches(0.6), size=16, color=LIGHT_BLUE)

# ── Slide 7: סיכום ──────────────────────────────────────────────────────────
s = add_slide(prs)
bg(s, DARK_BLUE)
box(s, 0, 0, W, Inches(0.12), ACCENT_BLUE)
box(s, 0, H - Inches(0.12), W, Inches(0.12), ACCENT_BLUE)

txt(s, "תודה 📐",
    Inches(0), Inches(1.6), W, Inches(1.4),
    size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "מערכת ניהול חכמה למשרד מדידות קרקע",
    Inches(1), Inches(3.1), Inches(11), Inches(0.8),
    size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

txt(s, "http://localhost:8501",
    Inches(1), Inches(4.1), Inches(11), Inches(0.7),
    size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

tags2 = ["Python", "Streamlit", "Pandas", "Plotly", "scikit-learn"]
tag_w2 = Inches(1.7)
start2 = Inches(2.6)
for i, tag in enumerate(tags2):
    box(s, start2 + i * (tag_w2 + Inches(0.15)),
        Inches(5.2), tag_w2, Inches(0.5), MID_BLUE)
    txt(s, tag,
        start2 + i * (tag_w2 + Inches(0.15)) + Inches(0.05),
        Inches(5.22), tag_w2 - Inches(0.1), Inches(0.5),
        size=14, bold=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────────────────
out = r"c:\Users\Ps\Desktop\סמסטר 22\ai\antigrafity\antigrafity_presentation.pptx"
prs.save(out)
print("Saved OK")
